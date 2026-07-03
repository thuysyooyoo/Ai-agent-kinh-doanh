# -*- coding: utf-8 -*-
import sys
import copy
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.stdout.reconfigure(encoding='utf-8')

# Colors
C_SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
C_SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
C_SLATE_600 = RGBColor(0x47, 0x55, 0x69)
C_SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
C_SLATE_400 = RGBColor(0x94, 0xA3, 0xB8)
C_SLATE_50  = RGBColor(0xF8, 0xFA, 0xFC)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

C_RED_600    = RGBColor(0xDC, 0x26, 0x26)
C_ORANGE_600 = RGBColor(0xEA, 0x58, 0x0C)
C_BLUE_600   = RGBColor(0x25, 0x63, 0xEB)
C_GREEN_600  = RGBColor(0x16, 0xA3, 0x4A)
C_SKY_600    = RGBColor(0x02, 0x84, 0xC7)
C_TEAL_600   = RGBColor(0x0D, 0x94, 0x88)
C_INDIGO_600 = RGBColor(0x4F, 0x46, 0xE5)
C_VIOLET_600 = RGBColor(0x7C, 0x3A, 0xED)

def duplicate_template_slide(prs, source_slide_idx, new_title, new_subtitle):
    """
    Creates a new slide based on a source slide's background, header, and footer.
    Deletes content elements (tables, callout shapes) so it's a clean template.
    """
    source_slide = prs.slides[source_slide_idx]
    new_slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    
    # Copy template elements only
    for shape in source_slide.shapes:
        if shape.has_table or "Table" in shape.name or "Rounded Rectangle 10" in shape.name or "TextBox 11" in shape.name:
            continue
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.append(new_el)
        
    # Update Title and Subtitle
    for shape in new_slide.shapes:
        if shape.name == "TextBox 2" and shape.has_text_frame:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = new_title
            p.font.name = 'Aptos'
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = C_SLATE_800
        elif shape.name == "TextBox 3" and shape.has_text_frame:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = new_subtitle
            p.font.name = 'Aptos'
            p.font.size = Pt(12)
            p.font.color.rgb = C_SLATE_500
            
    return new_slide

def move_slide(prs, slide, target_idx):
    """
    Moves a slide to target_idx (0-based index)
    """
    sldIdLst = prs.slides._sldIdLst
    slide_id = prs.slides.index(slide)
    slide_element = sldIdLst[slide_id]
    sldIdLst.remove(slide_element)
    sldIdLst.insert(target_idx, slide_element)

def add_table_to_slide(slide, rows, cols, left, top, width, height, col_widths, data, header_colors, font_size=8.5):
    """
    Adds a custom table, sets cell colors, column widths, and formats paragraphs.
    """
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set widths
    for idx, w in enumerate(col_widths):
        table.columns[idx].width = w
        
    # Populate cell content & styles
    for r_idx, row_data in enumerate(data):
        for c_idx, text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            
            # cell background
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = header_colors[c_idx]
            else:
                if r_idx % 2 == 1:
                    cell.fill.fore_color.rgb = C_SLATE_50
                else:
                    cell.fill.fore_color.rgb = C_WHITE
            
            # cell text
            cell.text_frame.clear()
            cell.text_frame.margin_top = Inches(0.08)
            cell.text_frame.margin_bottom = Inches(0.08)
            cell.text_frame.margin_left = Inches(0.08)
            cell.text_frame.margin_right = Inches(0.08)
            
            lines = text.split('\n')
            for line_idx, line in enumerate(lines):
                if line_idx > 0:
                    p = cell.text_frame.add_paragraph()
                else:
                    p = cell.text_frame.paragraphs[0]
                
                # bullet points or lists
                p.text = line
                p.font.name = 'Aptos'
                p.font.size = Pt(font_size)
                
                if r_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = C_WHITE
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = C_SLATE_800
                    if line.strip().startswith('•') or line.strip().startswith('-'):
                        p.alignment = PP_ALIGN.LEFT
                    else:
                        # left align for text columns, center for simple columns
                        if c_idx == 0 or len(row_data) <= 4 or '\n' in text:
                            p.alignment = PP_ALIGN.LEFT
                        else:
                            p.alignment = PP_ALIGN.CENTER
    return table_shape

def update_textbox_paragraphs(shape, paragraphs_text):
    """
    Updates the text of a textbox's paragraphs.
    """
    if not shape.has_text_frame:
        return
    for idx, text in enumerate(paragraphs_text):
        if idx < len(shape.text_frame.paragraphs):
            p = shape.text_frame.paragraphs[idx]
            p.text = text
        else:
            p = shape.text_frame.add_paragraph()
            p.text = text
        # maintain standard font styles
        p.font.name = 'Aptos'
        if idx == 0 and len(paragraphs_text) > 1:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = C_SLATE_800
        else:
            p.font.size = Pt(10)
            p.font.color.rgb = C_SLATE_600

def main():
    prs = Presentation(r"C:\Users\Admin\Downloads\Webinar_Hang_Rui_Ro_TXNG_04072026_FINAL_REDESIGNED_v3.pptx")
    print("Presentation loaded. Slides count:", len(prs.slides))
    
    # ----------------------------------------------------
    #  SLIDE 2: NỘI DUNG WEBINAR
    # ----------------------------------------------------
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.name == "TextBox 21" and shape.has_text_frame:
            # Sửa tiền kiểm / hậu kiểm
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = u"Hàng rủi ro cao: Quy trình 5 bước — kiểm tra trước thông quan. Hàng rủi ro trung bình: Quy trình 3 bước — tự công bố trước khi lưu thông. Làm rõ: đưa hàng về bảo quản không đồng nghĩa với đã thông quan."
            p.font.name = 'Aptos'
            p.font.size = Pt(10)
            p.font.color.rgb = C_SLATE_600
            
    # ----------------------------------------------------
    #  SLIDE 3: BỐI CẢNH PHÁP LÝ MỚI
    # ----------------------------------------------------
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.name == "TextBox 23" and shape.has_text_frame:
            shape.text_frame.clear()
            p0 = shape.text_frame.paragraphs[0]
            p0.text = u"Danh mục phân loại"
            p0.font.name = 'Aptos'
            p0.font.size = Pt(12)
            p0.font.bold = True
            p0.font.color.rgb = C_SLATE_800
            
            p1 = shape.text_frame.add_paragraph()
            p1.text = u"Các Bộ quản lý"
            p1.font.name = 'Aptos'
            p1.font.size = Pt(11)
            p1.font.color.rgb = C_SLATE_600
            
            p2 = shape.text_frame.add_paragraph()
            p2.text = u"• Bộ KH&CN: 135 sản phẩm (44 cao + 91 trung bình - Thông tư 36/2026/TT-BKHCN)."
            p2.font.name = 'Aptos'
            p2.font.size = Pt(10.5)
            p2.font.color.rgb = C_SLATE_600
            
            p3 = shape.text_frame.add_paragraph()
            p3.text = u"• Bộ Xây dựng (GTVT): 44 sản phẩm (Thông tư 49/2026/TT-BXD)."
            p3.font.name = 'Aptos'
            p3.font.size = Pt(10.5)
            p3.font.color.rgb = C_SLATE_600
            
            p4 = shape.text_frame.add_paragraph()
            p4.text = u"• Bộ NN&MT: 69 sản phẩm (nhiều hàng trung bình nhất)."
            p4.font.name = 'Aptos'
            p4.font.size = Pt(10.5)
            p4.font.color.rgb = C_SLATE_600
            
            p5 = shape.text_frame.add_paragraph()
            p5.text = u"• Dự kiến lần lượt ban hành 2026-2027."
            p5.font.name = 'Aptos'
            p5.font.size = Pt(10.5)
            p5.font.color.rgb = C_SLATE_600

    # ----------------------------------------------------
    #  SLIDE 4: 8 BỘ QUẢN LÝ TABLE
    # ----------------------------------------------------
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_table:
            table = shape.table
            # Row 1 (BKHCN)
            table.cell(1, 3).text = u"Thiết bị điện, đồ chơi trẻ em (Thông tư 36/2026/TT-BKHCN)"
            # Row 2 (BCT)
            table.cell(2, 3).text = u"VLN, thép, sữa, dệt may (Thông tư 33/2026/TT-BCT)"
            # Row 3 (BXD - GTVT)
            table.cell(3, 3).text = u"Ô tô, xe máy, phụ tùng xe (Thông tư 49/2026/TT-BXD)"
            # Row 4 (BXD - VLXD)
            table.cell(4, 3).text = u"Xi măng, kính, thạch cao (Thông tư 41/2026/TT-BXD)"
            # Row 6 (BYT)
            table.cell(6, 3).text = u"Dược liệu, TPCN, trang thiết bị y tế (Thông tư 24/2026/TT-BYT)"
            
            # format the modified cells
            for r_idx in [1, 2, 3, 4, 6]:
                cell = table.cell(r_idx, 3)
                p = cell.text_frame.paragraphs[0]
                p.font.name = 'Aptos'
                p.font.size = Pt(9.5)
                p.font.color.rgb = C_SLATE_800

    # ----------------------------------------------------
    #  SLIDE 6: HÀNG RỦI RO CAO PT7 (MODIFICATIONS)
    # ----------------------------------------------------
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.name == "TextBox 2" and shape.has_text_frame:
            shape.text_frame.paragraphs[0].text = u"🔴 Hàng Rủi ro Cao — Quy trình Đánh giá theo lô (Phương thức 7)"
        elif shape.name == "TextBox 13" and shape.has_text_frame:
            # Step 1
            update_textbox_paragraphs(shape, [
                u"Chuẩn bị trước nhập khẩu",
                u"Liên hệ và ký hợp đồng với Tổ chức chứng nhận chỉ định để sẵn sàng đánh giá lấy mẫu khi hàng về."
            ])
        elif shape.name == "TextBox 20" and shape.has_text_frame:
            # Step 2
            update_textbox_paragraphs(shape, [
                u"Hàng về cảng & Đăng ký kiểm tra",
                u"Đăng ký kiểm tra nhà nước trên Cổng một cửa quốc gia đính kèm hợp đồng chứng nhận → đưa hàng về kho bảo quản."
            ])
        elif shape.name == "TextBox 46" and shape.has_text_frame:
            # Callout
            update_textbox_paragraphs(shape, [
                u" Tuyệt đối: Hàng trong kho ở Bước 2-3 chưa được thông quan, chưa được bán. Chỉ được phép lưu thông sau khi có Thông báo đạt chất lượng và hoàn tất truy xuất nguồn gốc.",
                u"Trong thời gian chờ kết quả thử nghiệm: Hàng được lưu kho bảo quản — nghiêm cấm các hành vi bán hoặc đưa ra thị trường."
            ])

    # ----------------------------------------------------
    #  SLIDE 7: HÀNG RỦI RO TRUNG BÌNH (MODIFICATIONS)
    # ----------------------------------------------------
    slide7 = prs.slides[6]
    for shape in slide7.shapes:
        if shape.name == "TextBox 2" and shape.has_text_frame:
            shape.text_frame.paragraphs[0].text = u"🟡 Hàng Rủi ro Trung bình — Quy trình 3 bước (PT1, PT5 hoặc PT7)"
        elif shape.name == "TextBox 3" and shape.has_text_frame:
            shape.text_frame.paragraphs[0].text = u"Áp dụng cho 224 sản phẩm rủi ro trung bình | Tự công bố hợp quy | Thời gian: 1-3 ngày"
        elif shape.name == "TextBox 13" and shape.has_text_frame:
            # Step 1
            update_textbox_paragraphs(shape, [
                u"Đánh giá sự phù hợp",
                u"Chọn Phương thức 1 (Tự đánh giá mẫu điển hình), Phương thức 7 (Đánh giá theo lô) hoặc Phương thức 5 (Chứng nhận 03 năm) để đánh giá sản phẩm."
            ])
        elif shape.name == "TextBox 20" and shape.has_text_frame:
            # Step 2
            update_textbox_paragraphs(shape, [
                u"Nộp hồ sơ công bố",
                u"Đăng ký công bố hợp quy trên nqi.gov.vn. Hàng hóa nhập khẩu được truyền tờ khai và thông quan ngay, không cần kiểm tra chất lượng tại cảng."
            ])
        elif shape.name == "TextBox 27" and shape.has_text_frame:
            # Step 3
            update_textbox_paragraphs(shape, [
                u"Kiểm tra đột xuất & TXNG",
                u"Hàng hóa lưu thông chịu giám sát đột xuất. Khuyến nghị đăng ký GS1 và dán mã QR truy xuất nguồn gốc để hưởng chế độ giảm kiểm tra."
            ])

    # ----------------------------------------------------
    #  SLIDE 8: COMPARISON TABLE (MODIFICATIONS)
    # ----------------------------------------------------
    slide8 = prs.slides[7]
    for shape in slide8.shapes:
        if shape.has_table:
            table = shape.table
            
            # Headers
            table.cell(0, 0).text = u"Tiêu chí"
            table.cell(0, 1).text = u"🔴 RỦI RO CAO"
            table.cell(0, 2).text = u"🟡 RỦI RO TRUNG BÌNH"
            
            # Row 1
            table.cell(1, 0).text = u"Kiểm tra chất lượng"
            table.cell(1, 1).text = u"Kiểm tra chất lượng trước thông quan\n(Phương thức 7 theo lô hoặc Phương thức 5)"
            table.cell(1, 2).text = u"Tự đánh giá chất lượng sản phẩm\n(Phương thức 1, 5 hoặc 7)"
            
            # Row 2
            table.cell(2, 0).text = u"Công bố hợp quy"
            table.cell(2, 1).text = u"Được Tổ chức chứng nhận chỉ định cấp Giấy chứng nhận hợp quy"
            table.cell(2, 2).text = u"Doanh nghiệp tự công bố lên CSDL quốc gia (hoặc thuê Tổ chức chứng nhận)"
            
            # Row 3
            table.cell(3, 0).text = u"Thời điểm thông quan"
            table.cell(3, 1).text = u"Chỉ sau khi có thông báo đạt từ Cơ quan kiểm tra nhà nước"
            table.cell(3, 2).text = u"Thông quan ngay\nKhông cần chờ kết quả"
            
            # Row 4
            table.cell(4, 0).text = u"Được bán hàng?"
            table.cell(4, 1).text = u"Chỉ sau khi có Thông báo đạt và hoàn tất truy xuất nguồn gốc"
            table.cell(4, 2).text = u"Có thể bán ngay\nSau khi thông quan"
            
            # Row 5
            table.cell(5, 0).text = u"Truy xuất nguồn gốc"
            table.cell(5, 1).text = u"Bắt buộc định danh GS1 và dán mã QR theo lô"
            table.cell(5, 2).text = u"Khuyến khích thực hiện để hưởng ưu đãi"
            
            # Row 6
            table.cell(6, 0).text = u"Thời gian thông quan"
            table.cell(6, 1).text = u"7-14 ngày đối với PT7\n(hoặc 1-2 ngày đối với PT5 đã có GCN)"
            table.cell(6, 2).text = u"1-3 ngày\n(Thông quan ngay)"
            
            # Row 7
            table.cell(7, 0).text = u"Kiểm tra đột xuất"
            table.cell(7, 1).text = u"Có — Tần suất kiểm tra đột xuất cao hơn"
            table.cell(7, 2).text = u"Có — Tần suất thấp hơn nhưng vẫn chịu giám sát"
            
            # Format comparison table cells
            header_colors = [C_SLATE_900, C_RED_600, C_ORANGE_600]
            for r_idx in range(len(table.rows)):
                for c_idx in range(len(table.columns)):
                    cell = table.cell(r_idx, c_idx)
                    cell.fill.solid()
                    if r_idx == 0:
                        cell.fill.fore_color.rgb = header_colors[c_idx]
                    else:
                        if r_idx % 2 == 1:
                            cell.fill.fore_color.rgb = C_SLATE_50
                        else:
                            cell.fill.fore_color.rgb = C_WHITE
                    
                    p = cell.text_frame.paragraphs[0]
                    p.font.name = 'Aptos'
                    p.font.size = Pt(9.5)
                    if r_idx == 0:
                        p.font.bold = True
                        p.font.color.rgb = C_WHITE
                        p.alignment = PP_ALIGN.CENTER
                    else:
                        p.font.color.rgb = C_SLATE_800
                        p.alignment = PP_ALIGN.LEFT

    # ----------------------------------------------------
    #  SLIDE 12 & 13: SỬA HẬU KIỂM
    # ----------------------------------------------------
    slide12 = prs.slides[11]
    for shape in slide12.shapes:
        if shape.name == "Rounded Rectangle 20" and shape.has_text_frame:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = u"💡  Doanh nghiệp tự xác định & tự lưu hồ sơ: Tự đánh giá đủ điều kiện. Lưu 3 bộ thông báo đạt để phục vụ kiểm tra đột xuất — KHÔNG CẦN nộp đơn xin."
            p.font.name = 'Aptos'
            p.font.size = Pt(11)
            p.font.color.rgb = C_SLATE_800
        elif shape.name == "TextBox 29" and shape.has_text_frame:
            text = shape.text_frame.text.replace(u"mẫu hậu kiểm", u"mẫu kiểm tra đột xuất")
            shape.text_frame.clear()
            lines = text.split('\n')
            for l_idx, line in enumerate(lines):
                p = shape.text_frame.add_paragraph() if l_idx > 0 else shape.text_frame.paragraphs[0]
                p.text = line
                p.font.name = 'Aptos'
                p.font.size = Pt(10)
                p.font.color.rgb = C_SLATE_600

    slide13 = prs.slides[12]
    for shape in slide13.shapes:
        if shape.name == "TextBox 34" and shape.has_text_frame:
            # Step 4
            update_textbox_paragraphs(shape, [
                u"Kiểm tra đột xuất & Báo cáo định kỳ",
                u"Nộp báo cáo định kỳ 3 tháng/lần về tình hình nhập khẩu hàng hóa → Cơ quan chức năng tiến hành kiểm tra đột xuất."
            ])

    # ----------------------------------------------------
    #  CREATE NEW SLIDE 6B: PT5 HIGH-RISK PROCESS (INDEX 6)
    # ----------------------------------------------------
    # Slide 6 is index 5. We want the new slide right after it, making it index 6.
    slide6b = duplicate_template_slide(prs, 5, 
                                      u"🔴 Hàng Rủi ro Cao — Quy trình Chứng nhận Dài hạn (Phương thức 5)", 
                                      u"Áp dụng cho các sản phẩm nhập khẩu ổn định | Giấy chứng nhận 3 năm | Thời gian: 1-2 ngày")
    
    # Update step textboxes on Slide 6b
    for shape in slide6b.shapes:
        if shape.name == "TextBox 13" and shape.has_text_frame:
            update_textbox_paragraphs(shape, [
                u"Đánh giá xưởng & Thử mẫu",
                u"Đánh giá quy trình sản xuất của nhà máy nước ngoài + thử nghiệm mẫu điển hình."
            ])
        elif shape.name == "TextBox 20" and shape.has_text_frame:
            update_textbox_paragraphs(shape, [
                u"Cấp GCN hợp quy 3 năm",
                u"Được cấp Giấy chứng nhận hợp quy dài hạn (hiệu lực không quá 03 năm) cấp cho dòng sản phẩm."
            ])
        elif shape.name == "TextBox 27" and shape.has_text_frame:
            update_textbox_paragraphs(shape, [
                u"Đăng ký & Thông quan nhanh",
                u"Hàng về, đăng ký kiểm tra trên Cổng một cửa quốc gia đính kèm GCN → Cơ quan ra TB đạt trong 01 ngày → Thông quan."
            ])
        elif shape.name == "TextBox 34" and shape.has_text_frame:
            update_textbox_paragraphs(shape, [
                u"TXNG & Lưu thông",
                u"Dán nhãn hợp quy CR, nhãn phụ tiếng Việt và dán mã QR truy xuất nguồn gốc theo lô để lưu thông hợp pháp."
            ])
        elif shape.name == "TextBox 41" and shape.has_text_frame:
            update_textbox_paragraphs(shape, [
                u"Đánh giá giám sát định kỳ",
                u"Đánh giá giám sát định kỳ hàng năm (tần suất ít nhất 1 lần/năm) tại nhà máy để duy trì hiệu lực của GCNHQ 3 năm."
            ])
        elif shape.name == "TextBox 46" and shape.has_text_frame:
            update_textbox_paragraphs(shape, [
                u" Tuyệt đối: Hàng trong kho ở Bước 1-2 chưa được thông quan, chưa được bán. Chỉ được phép lưu thông sau khi có Thông báo đạt từ cơ quan kiểm tra.",
                u"Hàng hóa nhập khẩu theo Phương thức 5 được rút ngắn thời gian kiểm tra thực tế, giúp thông quan nhanh chóng."
            ])
            
    # Move slide 6b to index 6
    move_slide(prs, slide6b, 6)

    # ----------------------------------------------------
    #  CREATE NEW SLIDE 8B: 8 METHODS COMPARISON (INDEX 9)
    # ----------------------------------------------------
    # Slide 8 is index 7. The new index after re-ordering slide 6b:
    # Original slide 8 index is 7, now it is index 8 because of slide 6b.
    # We want Slide 8b at index 9.
    slide8b = duplicate_template_slide(prs, 8,
                                      u"Các Phương thức Đánh giá Sự phù hợp (PT1 đến PT8) theo TT 14/2026/TT-BKHCN",
                                      u"Doanh nghiệp đối chiếu để lựa chọn phương thức đánh giá phù hợp nhất cho sản phẩm")
    
    # Table data for 8 methods
    data_8b = [
        [u"Tiêu chí", u"PT 1", u"PT 2", u"PT 3", u"PT 4", u"PT 5", u"PT 6", u"PT 7", u"PT 8"],
        [u"Đánh giá nhà máy", u"Không", u"Có", u"Có", u"Có", u"Có (Kỹ)", u"Có (QMS)", u"Không", u"Không"],
        [u"Lấy mẫu giám sát", u"Không", u"Thị trường", u"Nơi SX", u"Cả hai", u"Nơi SX/Thị trường", u"Không", u"Không", u"Không"],
        [u"Thử nghiệm mẫu", u"Mẫu đại diện", u"Mẫu đại diện", u"Mẫu đại diện", u"Mẫu đại diện", u"Mẫu đại diện", u"Không", u"Theo lô", u"100% SP"],
        [u"Thời hạn hiệu lực", u"Max 1 năm", u"Max 3 năm", u"Max 3 năm", u"Max 3 năm", u"Max 3 năm", u"Max 3 năm", u"Theo lô", u"Theo SP"],
        [u"Bên đánh giá", u"Tự đánh giá/TCCN", u"Tổ chức CN", u"Tổ chức CN", u"Tổ chức CN", u"Tổ chức CN", u"Tổ chức CN", u"Tự đánh giá/TCCN", u"Tổ chức KĐ"],
        [u"Phù hợp nhất cho", u"Hàng nhập lẻ ổn định thấp", u"Hàng nội địa lưu thông lớn", u"Hàng SX trong nước đi đại lý", u"Hàng rủi ro cao, bán rộng rãi", u"Hàng Nhóm 2 nhập ổn định", u"Dịch vụ, phần mềm", u"Lô hàng nhập đơn chiếc", u"Thiết bị giá trị lớn (nồi hơi...)"]
    ]
    
    header_colors_8b = [C_SLATE_900, C_SKY_600, C_TEAL_600, C_GREEN_600, C_INDIGO_600, C_RED_600, C_VIOLET_600, C_ORANGE_600, C_BLUE_600]
    col_widths_8b = [Inches(1.73)] + [Inches(1.25)] * 8
    
    # Table dimensions: same top/left/width/height as Slide 8 Table
    add_table_to_slide(slide8b, 7, 9, Inches(0.80), Inches(1.51), Inches(11.73), Inches(3.90), col_widths_8b, data_8b, header_colors_8b, font_size=8.5)
    move_slide(prs, slide8b, 9)

    # ----------------------------------------------------
    #  CREATE NEW SLIDE 8C: DOSSIER COMPARISON (INDEX 10)
    # ----------------------------------------------------
    slide8c = duplicate_template_slide(prs, 8,
                                      u"Phân biệt: Hồ sơ nộp Nhà nước vs. Hồ sơ Kỹ thuật lưu trữ nội bộ",
                                      u"Doanh nghiệp thực hiện theo quy định tại Điều 13 và Điều 14 của Thông tư 14/2026/TT-BKHCN")
    
    data_8c = [
        [u"Tiêu chí", u"1. Hồ sơ Đăng ký công bố (Nộp nhà nước - Điều 13)", u"2. Hồ sơ Kỹ thuật lưu trữ (Lưu tại DN - Điều 14)"],
        [u"Mục đích", u"Để cơ quan nhà nước thẩm định và cấp Phiếu tiếp nhận bản công bố hợp quy.", u"Để trình khi có thanh tra chất lượng liên ngành hoặc Quản lý thị trường kiểm tra."],
        [u"Nơi nộp/lưu", u"Nộp trực tuyến qua cổng nqi.gov.vn hoặc dịch vụ công địa phương.", u"Lưu giữ tại văn phòng hoặc nhà máy của doanh nghiệp trong tối thiểu 10 năm."],
        [u"Chứng từ bắt buộc", u"• Bản đăng ký công bố (Phụ lục IV)\n• Báo cáo tự đánh giá (Phụ lục V)\n• Danh mục sản phẩm (Phụ lục 01)\n• Bản sao Đăng ký kinh doanh\n• Phiếu kết quả test / GCNHQ", u"• Toàn bộ hồ sơ ở cột 1 (kèm Phiếu tiếp nhận gốc)\n• Tài liệu thiết kế, bản vẽ kỹ thuật chi tiết\n• Chứng từ nhập khẩu lô hàng (Tờ khai, Invoice, PKL...)\n• Hồ sơ kiểm soát chất lượng sản xuất"],
        [u"Lưu ý bảo mật", u"Không cần nộp bản vẽ chi tiết hoặc bí mật công nghệ lên hệ thống dịch vụ công công cộng.", u"Bắt buộc phải có đầy đủ hồ sơ thiết kế, hướng dẫn sử dụng tiếng Việt để cơ quan đối chiếu khi hậu kiểm."]
    ]
    
    header_colors_8c = [C_SLATE_900, C_BLUE_600, C_GREEN_600]
    col_widths_8c = [Inches(1.73), Inches(5.0), Inches(5.0)]
    
    add_table_to_slide(slide8c, 5, 3, Inches(0.80), Inches(1.51), Inches(11.73), Inches(3.90), col_widths_8c, data_8c, header_colors_8c, font_size=9.2)
    move_slide(prs, slide8c, 10)

    # ----------------------------------------------------
    #  CREATE NEW SLIDE 8D: TIMELINE OF 2839 (INDEX 11)
    # ----------------------------------------------------
    slide8d = duplicate_template_slide(prs, 8,
                                      u"Thời hạn giải quyết thủ tục hành chính (Quyết định 2839/QĐ-BKHCN)",
                                      u"Quy định thời gian giải quyết tối đa đối với kiểm tra nhà nước và công bố hợp quy")
    
    data_8d = [
        [u"Thủ tục hành chính", u"Thời hạn giải quyết tối đa", u"Ghi chú / Yêu cầu bổ sung"],
        [u"Đăng ký kiểm tra nhà nước về chất lượng (Đủ hồ sơ)", u"01 ngày làm việc", u"Cơ quan ra Thông báo đạt chất lượng (Mẫu 03 Phụ lục VII NĐ 37) để thông quan lô hàng."],
        [u"Bổ sung hồ sơ kiểm tra còn thiếu", u"07 ngày làm việc", u"Quá hạn phải có văn bản gia hạn nêu rõ lý do và dự kiến ngày hoàn thành bổ sung."],
        [u"Khắc phục lỗi ghi nhãn sản phẩm", u"05 ngày làm việc", u"Lô hàng chỉ được thông quan khi có bằng chứng khắc phục nhãn đạt yêu cầu chất lượng."],
        [u"Xét duyệt và tiếp nhận hồ sơ công bố hợp quy", u"05 ngày làm việc", u"Cơ quan cấp Phiếu tiếp nhận công bố hợp quy; nếu hồ sơ lỗi thông báo trong 03 ngày."]
    ]
    
    header_colors_8d = [C_SLATE_900, C_ORANGE_600, C_SLATE_600]
    col_widths_8d = [Inches(3.23), Inches(3.0), Inches(5.5)]
    
    add_table_to_slide(slide8d, 5, 3, Inches(0.80), Inches(1.51), Inches(11.73), Inches(3.90), col_widths_8d, data_8d, header_colors_8d, font_size=9.2)
    move_slide(prs, slide8d, 11)

    # ----------------------------------------------------
    #  DYNAMICALLY UPDATE ALL SLIDE FOOTER PAGE NUMBERS
    # ----------------------------------------------------
    update_page_numbers(prs)
    
    # Save the updated presentation
    output_path = r"C:\Users\Admin\Downloads\Webinar_Hang_Rui_Ro_TXNG_04072026_FINAL_REDESIGNED_v3.pptx"
    prs.save(output_path)
    print("FINISHED UPDATING POWERPOINT PRESENTATION successfully!")
    print(f"Total slides in updated deck: {len(prs.slides)}")

def update_page_numbers(prs):
    total = len(prs.slides)
    for idx, slide in enumerate(prs.slides, 1):
        if idx == 1:
            continue
        for shape in slide.shapes:
            if shape.name == "TextBox 6" and shape.has_text_frame:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = f"{idx} / {total}"
                p.font.name = 'Aptos'
                p.font.size = Pt(9.5)
                p.font.color.rgb = C_SLATE_400
                p.alignment = PP_ALIGN.RIGHT

if __name__ == "__main__":
    main()
